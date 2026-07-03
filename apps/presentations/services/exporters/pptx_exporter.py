from __future__ import annotations

from io import BytesIO
from typing import Any, Sequence

from django.utils.text import slugify
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from apps.presentations.services.exporters import BaseExporter, ExportResult, register_exporter


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


_KIND_PREFIX = {"ok": "✓", "warn": "!", "risk": "✗", "num": "•", "info": "→"}


def _dicts(value) -> list:
    """Return only the dict items of a list; [] for anything malformed (so a
    dict/None where a list is expected never crashes the export)."""
    return [i for i in value if isinstance(i, dict)] if isinstance(value, list) else []


def _flatten_visual(visual: dict) -> list[str]:
    """Best-effort textual summary of a right-panel visual block."""
    if not isinstance(visual, dict):
        return []
    data = visual.get("data") or {}
    if not isinstance(data, dict):
        return []
    vtype = visual.get("type")
    out: list[str] = []
    if vtype == "dashboard":
        out = [f"{c.get('value', '')} {c.get('label', '')}".strip() for c in _dicts(data.get("cells"))]
    elif vtype in ("bar_chart", "comparison"):
        out = [f"{b.get('label', '')}: {b.get('display', b.get('value', ''))}".strip() for b in _dicts(data.get("bars"))]
    elif vtype == "card_list":
        out = [f"{c.get('title', '')} — {c.get('text', '')}".strip(" —") for c in _dicts(data.get("cards"))]
    elif vtype == "timeline":
        out = [f"{e.get('date', '')} {e.get('label', '')}".strip() for e in _dicts(data.get("events"))]
    elif vtype == "donut":
        out = [f"{data.get('center', data.get('percent', ''))} {data.get('label', '')}".strip()]
    elif vtype == "icon_grid":
        out = [c.get("label", "") for c in _dicts(data.get("cells"))]
    elif vtype == "status_card":
        out = [f"{data.get('title', '')} [{data.get('badge', '')}]".strip(" []"), data.get("text", "")]
    return [o for o in out if o]


def _flatten_content(content: dict) -> str:
    """Flatten a structured slide content dict into plain text for pptx export."""
    if not isinstance(content, dict):
        return ""
    lines: list[str] = []
    for key in ("eyebrow", "subtitle", "description"):
        if content.get(key):
            lines.append(content[key])
    for p in _dicts(content.get("points")):
        prefix = _KIND_PREFIX.get(p.get("kind", ""), "•")
        label, text = p.get("label", ""), p.get("text", "")
        body = f"{label} — {text}" if label and text else (label or text)
        lines.append(f"{prefix} {body}".strip())
    if content.get("highlight"):
        lines.append("» " + content["highlight"])
    lines.extend(_flatten_visual(content.get("visual") or {}))
    for s in _dicts(content.get("stats")):
        lines.append(f"{s.get('value', '')} {s.get('label', '')}".strip())
    if content.get("footer"):
        lines.append(content["footer"])
    return "\n".join(line for line in lines if line)


class PptxExporter(BaseExporter):
    def export(self, presentation: Any, slides: Sequence, theme: Any | None) -> ExportResult:
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        primary = _hex_to_rgb(theme.primary_color) if theme else RGBColor(0x1A, 0x36, 0x5D)
        secondary = _hex_to_rgb(theme.secondary_color) if theme else RGBColor(0xFF, 0xFF, 0xFF)
        accent = _hex_to_rgb(theme.accent_color) if theme else RGBColor(0x31, 0x82, 0xCE)
        font_heading = theme.font_heading if theme else "Arial"
        font_body = theme.font_body if theme else "Arial"

        for slide_data in slides:
            self._add_slide(prs, slide_data, primary, secondary, font_heading, font_body)

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        filename = f"{slugify(presentation.title) or 'presentation'}.pptx"
        return ExportResult(
            content=buffer.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=filename,
        )

    def _add_slide(
        self,
        prs: PptxPresentation,
        slide_data: Any,
        primary: RGBColor,
        secondary: RGBColor,
        font_heading: str,
        font_body: str,
    ) -> None:
        layout = prs.slide_layouts[6]  # blank layout
        pptx_slide = prs.slides.add_slide(layout)

        bg = pptx_slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = primary

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        is_title = getattr(slide_data, "slide_type", "split") in ("cover", "closing")

        if is_title:
            heading_top = Inches(2.5)
            heading_height = Inches(1.5)
            body_top = Inches(4.0)
        else:
            heading_top = Inches(0.7)
            heading_height = Inches(1.0)
            body_top = Inches(1.8)

        heading_box = pptx_slide.shapes.add_textbox(
            Inches(1), heading_top, slide_width - Inches(2), heading_height
        )
        heading_tf = heading_box.text_frame
        heading_tf.word_wrap = True
        heading_p = heading_tf.paragraphs[0]
        heading_p.text = slide_data.heading or ""
        heading_p.font.size = Pt(36) if is_title else Pt(28)
        heading_p.font.bold = True
        heading_p.font.color.rgb = secondary
        heading_p.font.name = font_heading
        if is_title:
            heading_p.alignment = 1  # CENTER

        body_box = pptx_slide.shapes.add_textbox(
            Inches(1), body_top, slide_width - Inches(2), slide_height - body_top - Inches(1)
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_text = _flatten_content(getattr(slide_data, "content", {}) or {})
        for i, line in enumerate(body_text.split("\n")):
            if i == 0:
                body_tf.paragraphs[0].text = line
                p = body_tf.paragraphs[0]
            else:
                p = body_tf.add_paragraph()
                p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = secondary
            p.font.name = font_body
            if is_title:
                p.alignment = 1

        notes = getattr(slide_data, "notes", "")
        if notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = notes


register_exporter("pptx", PptxExporter)
